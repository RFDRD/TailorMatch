"""
TailorMatch - Prezentatsiya Slaydlari Generator
Bu kod Google Colab'da ishlatiladi va PPTX prezentatsiya yaratadi
Redmi 14 Pro uchun screenshot joylari bilan
"""

# Kerakli kutubxonalarni o'rnatish
!pip install python-pptx Pillow -q

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from google.colab import files
from datetime import datetime

# Prezentatsiya yaratish (16:9 format)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Ranglar
PRIMARY_COLOR = RGBColor(233, 30, 99)  # Pink
SECONDARY_COLOR = RGBColor(156, 39, 176)  # Purple
DARK_BG = RGBColor(18, 18, 18)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
DARK_GRAY = RGBColor(66, 66, 66)

def add_gradient_background(slide, color1=DARK_BG, color2=PRIMARY_COLOR):
    """Gradient fon qo'shish"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color1
    background.line.fill.background()
    # Z-indexni pastga tushirish
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_phone_placeholder(slide, left, top, width=Inches(2.2), height=Inches(4.5), label="Screenshot"):
    """Telefon screenshot uchun joy"""
    # Telefon ramkasi
    phone_frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    phone_frame.fill.solid()
    phone_frame.fill.fore_color.rgb = RGBColor(40, 40, 40)
    phone_frame.line.color.rgb = RGBColor(80, 80, 80)
    phone_frame.line.width = Pt(2)
    
    # Ekran joyi
    screen = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.1),
        top + Inches(0.15),
        width - Inches(0.2),
        height - Inches(0.3)
    )
    screen.fill.solid()
    screen.fill.fore_color.rgb = LIGHT_GRAY
    screen.line.fill.background()
    
    # Label
    label_box = slide.shapes.add_textbox(
        left, top + height + Inches(0.1), width, Inches(0.4)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_title_slide(title, subtitle=""):
    """Sarlavha slayd"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)
    
    # Logo/Icon
    logo_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5)
    )
    logo_shape.fill.solid()
    logo_shape.fill.fore_color.rgb = PRIMARY_COLOR
    logo_shape.line.fill.background()
    
    # Logo text
    logo_text = slide.shapes.add_textbox(Inches(5.9), Inches(1.7), Inches(1.5), Inches(1.1))
    tf = logo_text.text_frame
    p = tf.paragraphs[0]
    p.text = "TM"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, bullets=None, show_phones=0, phone_labels=None):
    """Kontent slayd"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, DARK_BG)
    
    # Dekorativ chiziq
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(1), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_COLOR
    line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content width based on phones
    content_width = Inches(7) if show_phones > 0 else Inches(12)
    
    # Bullets
    if bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), content_width, Inches(5))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(220, 220, 220)
            p.space_after = Pt(14)
    
    # Phone placeholders
    if show_phones > 0:
        if phone_labels is None:
            phone_labels = [f"Screenshot {i+1}" for i in range(show_phones)]
        
        phone_spacing = Inches(2.6)
        start_x = Inches(7.8)
        
        for i in range(show_phones):
            if i < len(phone_labels):
                add_phone_placeholder(
                    slide,
                    start_x + (i * phone_spacing),
                    Inches(1.3),
                    label=phone_labels[i]
                )
    
    return slide

def add_features_slide(title, features, phone_count=2, phone_labels=None):
    """Xususiyatlar slayd"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_BG, DARK_BG)
    
    # Title with accent
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(0.8), Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY_COLOR
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Features as cards
    card_y = Inches(1.5)
    for i, feature in enumerate(features[:6]):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            card_y + (i * Inches(0.85)),
            Inches(6.5),
            Inches(0.75)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(35, 35, 35)
        card.line.color.rgb = RGBColor(60, 60, 60)
        
        # Icon circle
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.7),
            card_y + (i * Inches(0.85)) + Inches(0.15),
            Inches(0.45),
            Inches(0.45)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = PRIMARY_COLOR
        icon.line.fill.background()
        
        # Feature text
        text_box = slide.shapes.add_textbox(
            Inches(1.3),
            card_y + (i * Inches(0.85)) + Inches(0.2),
            Inches(5.5),
            Inches(0.5)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = feature
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(230, 230, 230)
    
    # Phone placeholders
    if phone_labels is None:
        phone_labels = [f"Screenshot {i+1}" for i in range(phone_count)]
    
    phone_spacing = Inches(2.6)
    start_x = Inches(7.8)
    
    for i in range(phone_count):
        if i < len(phone_labels):
            add_phone_placeholder(
                slide,
                start_x + (i * phone_spacing),
                Inches(1.3),
                label=phone_labels[i]
            )
    
    return slide

# ==================== SLAYDLAR ====================

# SLAYD 1: Titul
add_title_slide(
    "TailorMatch",
    "Chevarlar va mijozlar uchun mobil ilova • O'zbekiston"
)

# SLAYD 2: Muammo va yechim
add_content_slide(
    "Muammo va Yechim",
    [
        "Mijozlar yaxshi chevar topishda qiyinchiliklarga duch keladi",
        "Chevarlar yangi mijozlar topishda muammolar bor",
        "Buyurtma berish jarayoni noqulay va uzoq",
        "Aloqa usullari eskirgan (faqat telefon orqali)",
        "",
        "TailorMatch — bularning barchasiga zamonaviy yechim!"
    ]
)

# SLAYD 3: Ilova haqida
add_content_slide(
    "TailorMatch nima?",
    [
        "Flutter texnologiyasida ishlab chiqilgan mobil ilova",
        "Android, iOS va Web platformalarida ishlaydi",
        "Firebase backend bilan real-time sinxronizatsiya",
        "O'zbekcha, Ruscha, Inglizcha tillarni qo'llab-quvvatlaydi",
        "Foydalanish uchun bepul"
    ],
    show_phones=1,
    phone_labels=["Login oynasi"]
)

# SLAYD 4: Mijoz funksiyalari
add_features_slide(
    "Mijoz imkoniyatlari",
    [
        "Chevarlarni qidirish va ko'rish",
        "Chevar xizmatlari va narxlarni ko'rish",
        "Band kunlarni oldindan bilish",
        "Buyurtma berish",
        "Chat orqali muloqot qilish",
        "Baho berish (5 aspektli Yandex uslubida)"
    ],
    phone_count=2,
    phone_labels=["Asosiy ekran", "Chevar profili"]
)

# SLAYD 5: Chevar funksiyalari
add_features_slide(
    "Chevar imkoniyatlari",
    [
        "Buyurtmalarni boshqarish (4 kategoriya)",
        "Xizmatlarni qo'shish va tahrirlash",
        "Kalendarda band kunlarni belgilash",
        "Mijozlar bilan chat muloqoti",
        "Profil va sozlamalarni tahrirlash",
        "Ko'p tilli interfeys"
    ],
    phone_count=2,
    phone_labels=["Buyurtmalar", "Xizmatlar"]
)

# SLAYD 6: Chat tizimi
add_content_slide(
    "Real-time Chat Tizimi",
    [
        "Mijoz va chevar o'rtasida tezkor aloqa",
        "Real-time xabar almashish",
        "Xabar jo'natish/olish bildirishi",
        "Suhbat tarixini saqlash",
        "Qulay va zamonaviy interfeys"
    ],
    show_phones=2,
    phone_labels=["Mijoz chati", "Chevar chati"]
)

# SLAYD 7: Baho tizimi
add_content_slide(
    "Ko'p aspektli baho tizimi",
    [
        "Yandex Go uslubida 5 aspektli baholash:",
        "",
        "⭐ Sifat — tikilish sifati",
        "📏 Hajm — to'g'ri o'lcham",
        "⚡ Tezlik — buyurtma muddati",
        "🤝 Muomala — xizmat ko'rsatish",
        "💰 Narx — narx/sifat nisbati"
    ],
    show_phones=1,
    phone_labels=["Baho berish"]
)

# SLAYD 8: Texnologiyalar
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, DARK_BG, DARK_BG)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Texnologiya Steki"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Accent line
accent = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1), Inches(0.06)
)
accent.fill.solid()
accent.fill.fore_color.rgb = PRIMARY_COLOR
accent.line.fill.background()

# Tech cards
techs = [
    ("Flutter", "Cross-platform UI framework", Inches(0.5)),
    ("Dart", "Dasturlash tili", Inches(3.5)),
    ("Firebase", "Backend xizmatlari", Inches(6.5)),
    ("Firestore", "Real-time database", Inches(9.5)),
]

for tech, desc, x_pos in techs:
    # Card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.8), Inches(2.8), Inches(2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(40, 40, 40)
    card.line.color.rgb = PRIMARY_COLOR
    card.line.width = Pt(2)
    
    # Tech name
    name_box = slide.shapes.add_textbox(x_pos, Inches(2.1), Inches(2.8), Inches(0.6))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = tech
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(x_pos, Inches(2.7), Inches(2.8), Inches(0.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER

# Additional features
features_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12), Inches(2.5))
tf = features_box.text_frame
tf.word_wrap = True

features_text = [
    "✓ Google Fonts (Lato) — zamonaviy tipografiya",
    "✓ Material Design 3 — UI standartlari",
    "✓ Cached Network Image — optimallashtirish",
    "✓ Provider — state management"
]

for i, feat in enumerate(features_text):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = feat
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.space_after = Pt(8)

# SLAYD 9: Kelajak rejalari
add_content_slide(
    "Kelajak Rejalari",
    [
        "📱 Push bildirishnomalar tizimi",
        "📸 Portfolio (chevar ishlari galereyasi)",
        "💳 Onlayn to'lov integratsiyasi (Payme, Click)",
        "🗺️ Geolokatsiya — yaqin chevarlarni topish",
        "⭐ Premium obuna tizimi",
        "📊 Chevarlar uchun analitika paneli"
    ]
)

# SLAYD 10: Yakuniy
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

# Big circle
circle = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(5.4), Inches(1.5), Inches(2.5), Inches(2.5)
)
circle.fill.solid()
circle.fill.fore_color.rgb = PRIMARY_COLOR
circle.line.fill.background()

# Logo text
logo_text = slide.shapes.add_textbox(Inches(5.4), Inches(2), Inches(2.5), Inches(1.5))
tf = logo_text.text_frame
p = tf.paragraphs[0]
p.text = "TM"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Thank you
thanks = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(1))
tf = thanks.text_frame
p = tf.paragraphs[0]
p.text = "Rahmat!"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Contact info
contact = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1))
tf = contact.text_frame
p = tf.paragraphs[0]
p.text = "Savollar uchun tayyor"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(180, 180, 180)
p.alignment = PP_ALIGN.CENTER

# Date
date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
tf = date_box.text_frame
p = tf.paragraphs[0]
p.text = datetime.now().strftime("%d.%m.%Y")
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(120, 120, 120)
p.alignment = PP_ALIGN.CENTER

# Faylni saqlash va yuklab olish
filename = "TailorMatch_Prezentatsiya.pptx"
prs.save(filename)
files.download(filename)

print(f"✅ Prezentatsiya muvaffaqiyatli yaratildi: {filename}")
print("📥 Yuklab olish avtomatik boshlanadi...")
print()
print("=" * 60)
print("📸 SCREENSHOT QO'LLANMASI:")
print("=" * 60)
print()
print("Slayd 3 - Login oynasi:")
print("  → Login/Register ekranini screenshotini qo'ying")
print()
print("Slayd 4 - Mijoz imkoniyatlari:")
print("  → Asosiy ekran: Chevarlar ro'yxati (bosh sahifa)")
print("  → Chevar profili: Chevar tafsilotlari sahifasi")
print()
print("Slayd 5 - Chevar imkoniyatlari:")
print("  → Buyurtmalar: 4 tab'li buyurtmalar sahifasi")
print("  → Xizmatlar: Chevar xizmatlari ro'yxati")
print()
print("Slayd 6 - Chat tizimi:")
print("  → Mijoz chati: Mijoz tomonidan chat oynasi")
print("  → Chevar chati: Chevar tomonidan chat oynasi")
print()
print("Slayd 7 - Baho tizimi:")
print("  → Baho berish: 5 yulduzli + aspekt chiplar dialogi")
print()
print("=" * 60)
